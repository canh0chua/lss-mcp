import os
import tempfile
import ast
from pathlib import Path
import subprocess
import io
import base64
import json
import sqlite3
import shlex
import threading
import ipaddress
import logging
import requests
import re
import pathspec
import glob
from functools import lru_cache
from urllib.parse import urlparse
from mcp.server.fastmcp import FastMCP
import httpx
import tempfile

# --- Heavy libraries: lazy-loaded inside functions that need them ---
# fitz (PyMuPDF), docx, Presentation (python-pptx), openpyxl,
# BeautifulSoup, html2text, Image (PIL), pytesseract
# are imported on first use to keep baseline RSS low.

mcp = FastMCP("Local_AI_Support_Stack")

# ---------------------------------------------------------------------------
# Security helpers
# ---------------------------------------------------------------------------

# Allowed root directory for all local file access.
# Set WORKSPACE env var to override; falls back to ~/workspace for cross-platform local dev.
_WORKSPACE = os.path.realpath(os.environ.get("WORKSPACE", str(Path.home() / "workspace")))

# Optional project-level restriction within _WORKSPACE.
# When set, file tools are further restricted to this subdirectory.
# Example: WORKSPACE_PROJECT=my-repo  →  only /workspace/my-repo/** is accessible.
_WORKSPACE_PROJECT: str | None = None
_raw_project = os.environ.get("WORKSPACE_PROJECT", "").strip()
if _raw_project:
    _candidate = os.path.realpath(os.path.join(_WORKSPACE, _raw_project))
    if _candidate.startswith(_WORKSPACE + os.sep) or _candidate == _WORKSPACE:
        _WORKSPACE_PROJECT = _candidate
    else:
        print(f"Warning: WORKSPACE_PROJECT '{_raw_project}' escapes workspace — ignored.")

# Sensitive file patterns that should never be read or indexed.
_SENSITIVE_PATTERNS = (
    ".env", ".pem", ".key", ".p12", ".pfx", ".crt", ".cer",
    "id_rsa", "id_dsa", "id_ecdsa", "id_ed25519",
    ".netrc", ".pgp", ".gpg", "credentials", "secret",
)

# Extension allowlist for FTS indexing (text-based source files only).
_INDEXABLE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".rs", ".java",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".rb", ".php", ".swift",
    ".kt", ".scala", ".sh", ".bash", ".zsh", ".fish",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".json", ".xml", ".html", ".css", ".scss", ".md", ".txt",
    ".sql", ".graphql", ".proto", ".tf", ".hcl",
}

# Maximum input lengths for tool parameters.
_MAX_QUERY_LEN = 500
_MAX_PATH_LEN = 1000
_MAX_CMD_LEN = 2000
_MAX_URL_LEN = 2000

# Private/loopback IP ranges blocked for SSRF protection.
_PRIVATE_NETWORKS = [
    ipaddress.ip_network("10.0.0.0/8"),
    ipaddress.ip_network("172.16.0.0/12"),
    ipaddress.ip_network("192.168.0.0/16"),
    ipaddress.ip_network("127.0.0.0/8"),
    ipaddress.ip_network("169.254.0.0/16"),   # link-local / AWS metadata
    ipaddress.ip_network("::1/128"),
    ipaddress.ip_network("fc00::/7"),
]


def _validate_local_path(path: str) -> str:
    """
    Resolve *path* and assert it lives inside the allowed root.
    If WORKSPACE_PROJECT is set, the root is that subdirectory; otherwise _WORKSPACE.
    Returns the resolved absolute path on success.
    Raises ValueError with a safe message on failure.
    """
    if len(path) > _MAX_PATH_LEN:
        raise ValueError("Path exceeds maximum allowed length.")
    resolved = os.path.realpath(path)
    root = _WORKSPACE_PROJECT if _WORKSPACE_PROJECT else _WORKSPACE
    if not resolved.startswith(root + os.sep) and resolved != root:
        if _WORKSPACE_PROJECT:
            raise ValueError(
                f"Access denied: path is outside the configured project ({_WORKSPACE_PROJECT}). "
                f"IMPORTANT: Tell the user that WORKSPACE_PROJECT is not configured for this path. "
                f"They must set WORKSPACE_PROJECT in their .env file and restart the container. "
                f"Do NOT use this MCP for files outside the configured project — use native tools instead."
            )
        raise ValueError(
            f"Access denied: path is outside the allowed workspace ({_WORKSPACE})."
        )
    return resolved


def _is_sensitive_file(path: str) -> bool:
    """Return True if the filename matches a known-sensitive pattern."""
    name = os.path.basename(path).lower()
    return any(pat in name for pat in _SENSITIVE_PATTERNS)


def _validate_url(url: str, *, require_public: bool = True) -> str:
    """
    Validate *url* is a well-formed http/https URL.
    If *require_public* is True, reject URLs that resolve to private/loopback IPs.
    Returns the validated URL on success, raises ValueError on failure.
    """
    if len(url) > _MAX_URL_LEN:
        raise ValueError("URL exceeds maximum allowed length.")
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError(f"Unsupported URL scheme '{parsed.scheme}'. Only http/https are allowed.")
    if not parsed.netloc:
        raise ValueError("URL is missing a host.")
    if require_public:
        hostname = parsed.hostname or ""
        try:
            addr = ipaddress.ip_address(hostname)
            for net in _PRIVATE_NETWORKS:
                if addr in net:
                    raise ValueError(f"Access denied: URL resolves to a private/internal address.")
        except ValueError as exc:
            # Re-raise our own errors; ignore non-IP hostnames (DNS not resolved here).
            if "Access denied" in str(exc) or "Unsupported" in str(exc) or "missing" in str(exc):
                raise
    return url


def _validate_searxng_url(url: str) -> str:
    """Validate the SearXNG base URL (localhost is explicitly allowed)."""
    if len(url) > _MAX_URL_LEN:
        raise ValueError("SEARXNG_URL exceeds maximum allowed length.")
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError(f"SEARXNG_URL has unsupported scheme '{parsed.scheme}'.")
    if not parsed.netloc:
        raise ValueError("SEARXNG_URL is missing a host.")
    return url
def _validate_fourget_url(url: str) -> str:
    """Validate the 4get base URL (localhost is explicitly allowed)."""
    if len(url) > _MAX_URL_LEN:
        raise ValueError("FOURGET_URL exceeds maximum allowed length.")
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError(f"FOURGET_URL has unsupported scheme '{parsed.scheme}'.")
    if not parsed.netloc:
        raise ValueError("FOURGET_URL is missing a host.")
    return url

# ---------------------------------------------------------------------------
# Tools
# ---------------------------------------------------------------------------

@mcp.tool()
def web_search(query: str, type: str = "web", limit: int = 10, npt: str = "", scraper: str = "", nsfw: bool = False, country: str = "", lang: str = "", time_min: int = 0, time_max: int = 0) -> str:
    """Search the web using 4get (privacy-respecting proxy). Returns JSON results.

    Supported types:
      - web: general web results (title, url, snippet)
      - image: images (title, url, thumbnail, width, height)
      - video: videos (title, url, description, duration, views)
      - news: recent news articles (title, url, description, date, source)
      - music: music tracks (title, artist, album, stream url)

    Use 'npt' (next page token) from a previous response to paginate.
    """
    if len(query) > _MAX_QUERY_LEN:
        return f"Error: query exceeds maximum length of {_MAX_QUERY_LEN} characters."
    try:
        raw_url = os.getenv("FOURGET_URL", "http://localhost:8081")
        fourget_url = _validate_fourget_url(raw_url)
    except ValueError as e:
        return f"Configuration error: {e}"

    # Map type to 4get API endpoint and result key
    endpoint_map = {
        "web":    ("/api/v1/web",    "web"),
        "image":  ("/api/v1/images", "image"),
        "video":  ("/api/v1/videos", "video"),
        "news":   ("/api/v1/news",   "news"),
        "music":  ("/api/v1/music",  "song"),
    }

    if type not in endpoint_map:
        return f"Invalid type '{type}'. Supported: {', '.join(endpoint_map.keys())}"

    endpoint, result_key = endpoint_map[type]
    limit = max(1, min(limit, 20))

    try:
        params = {}
        if npt:
            params["npt"] = npt
        else:
            params["s"] = query
        if scraper:
            params["scraper"] = scraper
        if nsfw:
            params["nsfw"] = "1"
        if country:
            params["country"] = country
        if lang:
            params["lang"] = lang
        if time_min:
            params["time_min"] = time_min
        if time_max:
            params["time_max"] = time_max
        if country:
            params["country"] = country
        if lang:
            params["lang"] = lang
        if time_min:
            params["time_min"] = time_min
        if time_max:
            params["time_max"] = time_max
        if scraper:
            params["scraper"] = scraper
        if nsfw:
            params["nsfw"] = "1"
        if country:
            params["country"] = country
        if lang:
            params["lang"] = lang
        if time_min:
            params["time_min"] = time_min
        if time_max:
            params["time_max"] = time_max
        if country:
            params["country"] = country
        if lang:
            params["lang"] = lang
        if time_min:
            params["time_min"] = time_min
        if time_max:
            params["time_max"] = time_max
        if npt:
            params["npt"] = npt
        else:
            params["s"] = query

        resp = requests.get(
            f"{fourget_url}{endpoint}",
            params=params,
            timeout=15
        )
        data = resp.json()
        if data.get("status") != "ok":
            return f"Search failed: 4get returned status '{data.get('status', 'unknown')}'"

        # Build compact output per type
        if type == "web":
            raw = data.get(result_key, [])[:limit]
            out = [{"title": r.get("title", ""), "url": r.get("url", ""), "snippet": r.get("description", "")} for r in raw]
        elif type == "image":
            raw = data.get(result_key, [])[:limit]
            out = []
            for r in raw:
                sources = r.get("source", [])
                thumb = sources[-1].get("url", "") if sources else ""
                full = sources[0].get("url", "") if sources else ""
                out.append({"title": r.get("title", ""), "url": full, "thumbnail": thumb})
        elif type == "video":
            raw = data.get(result_key, [])[:limit]
            out = [{"title": r.get("title", ""), "url": r.get("url", ""), "description": r.get("description", ""), "duration": r.get("time", ""), "views": r.get("views", "")} for r in raw]
        elif type == "news":
            raw = data.get(result_key, [])[:limit]
            out = [{"title": r.get("title", ""), "url": r.get("url", ""), "description": r.get("description", ""), "date": r.get("date"), "source": r.get("source", "")} for r in raw]
        elif type == "music":
            raw = data.get(result_key, [])[:limit]
            out = [{"title": r.get("title", ""), "artist": r.get("artist", ""), "album": r.get("album", ""), "duration": r.get("time", "")} for r in raw]

        # Include spelling correction if available (web only)
        result = {"results": out}
        if type == "web":
            spelling = data.get("spelling", {})
            if spelling.get("type") != "no_correction":
                result["spelling"] = spelling
        # Include next page token if available
        if data.get("npt"):
            result["npt"] = data["npt"]

        return json.dumps(result)
    except Exception as e:
        return f"Search failed. Ensure 4get is running. Error: {str(e)}"


@mcp.tool()
async def web_extract_crw(url: str) -> str:
    """Fetch a URL, execute JavaScript, strip HTML bloat, and return pure Markdown.
    
    Uses CRW (Firecrawl-compatible API) for JS rendering and markdown extraction.
    No LLM involved — pure browser rendering + DOM-to-markdown.
    Requires the 'crw' and 'lightpanda' services to be running.
    """
    if len(url) > _MAX_URL_LEN:
        return f"Error: URL exceeds maximum length of {_MAX_URL_LEN} characters."
    try:
        crw_url = os.getenv("CRW_URL", "http://crw:3000")
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                f"{crw_url}/v1/scrape",
                json={"url": url, "formats": ["markdown"]},
                headers={"User-Agent": "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36"},
            )
            data = resp.json()
            if resp.status_code == 200 and data.get("success"):
                return data["data"]["markdown"]
            error_msg = data.get("error", str(data))
            return f"Failed to extract page content: {error_msg}"
    except Exception as e:
        return f"Search failed: {str(e)}"


@mcp.tool()
async def web_crawl(url: str, limit: int = 10) -> str:
    """Crawl a website starting from a URL. Returns Markdown for up to `limit` pages.

    Useful for mirroring a small site or reading multiple pages from a docs site.
    Returns a JSON array of {url, markdown} objects.
    """
    if len(url) > _MAX_URL_LEN:
        return f"Error: URL exceeds maximum length of {_MAX_URL_LEN} characters."
    try:
        crw_url = os.getenv("CRW_URL", "http://crw:3000")
        async with httpx.AsyncClient(timeout=120) as client:
            # Start crawl
            resp = await client.post(
                f"{crw_url}/v1/crawl",
                json={"url": url, "limit": limit, "scrapeOptions": {"formats": ["markdown"]}},
            )
            data = resp.json()
            if resp.status_code != 200 or not data.get("success"):
                return f"Crawl failed: {data.get('error', str(data))}"
            crawl_id = data.get("id")
            if not crawl_id:
                return f"Crawl failed: no crawl ID returned. Response: {str(data)}"

            # Poll until complete (max 60s)
            for _ in range(30):
                poll_resp = await client.get(f"{crw_url}/v1/crawl/{crawl_id}")
                poll_data = poll_resp.json()
                status = poll_data.get("status", "")
                if status == "completed":
                    pages = poll_data.get("data", [])
                    clean = [{"url": p.get("metadata", {}).get("sourceURL", ""), "markdown": p.get("markdown", "")} for p in pages]
                    return json.dumps(clean, indent=2)[:8000]
                if status == "failed":
                    return f"Crawl failed: {poll_data.get('error', 'unknown')}"
                import asyncio
                await asyncio.sleep(2)
            return f"Crawl timed out after 60s. ID: {crawl_id}"
    except httpx.RequestError as e:
        return f"Crawl failed: HTTP error connecting to CRW: {str(e)}"
    except Exception as e:
        return f"Crawl failed: {str(e)}"


@mcp.tool()
async def web_map(url: str, limit: int = 50) -> str:
    """Discover all pages on a website. Returns a list of URLs found on the site.

    Useful for finding documentation pages, blog posts, or sitemap entries.
    Returns a JSON array of {title, url} objects.
    """
    if len(url) > _MAX_URL_LEN:
        return f"Error: URL exceeds maximum length of {_MAX_URL_LEN} characters."
    try:
        crw_url = os.getenv("CRW_URL", "http://crw:3000")
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{crw_url}/v1/map",
                json={"url": url, "limit": limit},
            )
            data = resp.json()
            if resp.status_code == 200 and data.get("success"):
                links = data.get("data", {}).get("links", [])[:limit]
                clean = [{"title": "", "url": l} for l in links]
                return json.dumps(clean, indent=2)[:8000]
            return f"Map failed: {data.get('error', str(data))}"
    except httpx.RequestError as e:
        return f"Map failed: HTTP error connecting to CRW: {str(e)}"
    except Exception as e:
        return f"Map failed: {str(e)}"



def read_document(path_or_url: str) -> str:
    """
    Parse documents into Markdown using lightweight libraries.
    Supports both local files and remote URLs.

    Supported formats:
    - PDF (text only)
    - Word (.docx), PowerPoint (.pptx), Excel (.xlsx)
    - Images: PNG, JPEG, TIFF, BMP, GIF, WEBP (OCR via Tesseract)
    - HTML, XML
    - Plain text (.txt), CSV (.csv) converted to markdown tables

    Local paths are restricted to the workspace directory.
    Remote URLs must point to public hosts (private/internal IPs are blocked).

    Returns Markdown output. Caching is enabled for recently converted documents.
    """
    if len(path_or_url) > _MAX_URL_LEN:
        return f"Error: path/URL exceeds maximum allowed length."

    parsed = urlparse(path_or_url)
    is_url = parsed.scheme in ("http", "https")

    try:
        if is_url:
            try:
                _validate_url(path_or_url, require_public=True)
            except ValueError as e:
                return f"Error: {e}"
            source = path_or_url
        else:
            try:
                resolved = _validate_local_path(path_or_url)
            except ValueError as e:
                return f"Error: {e}"
            if _is_sensitive_file(resolved):
                return "Error: Access denied — sensitive file type."
            if not os.path.exists(resolved):
                return f"Error: Local file '{resolved}' not found."
            source = resolved

        return _convert_document_cached(source)
    except Exception as e:
        return f"Failed to parse document: {str(e)}"


@lru_cache(maxsize=20)
def _convert_document_cached(source: str) -> str:
    """Cache wrapper for document conversion."""
    return _convert_document(source)


def _convert_document(source: str) -> str:
    """Convert a document source (file path or URL) to Markdown."""
    if source.startswith(('http://', 'https://')):
        try:
            resp = requests.get(source, timeout=30)
            if resp.status_code != 200:
                return f"Failed to fetch URL: HTTP {resp.status_code}"
            content = resp.content
            content_type = resp.headers.get('Content-Type', '').split(';')[0]
            mime_map = {
                'application/pdf': '.pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                'image/jpeg': '.jpg', 'image/png': '.png', 'image/gif': '.gif',
                'image/webp': '.webp', 'text/html': '.html', 'application/xml': '.xml',
                'text/plain': '.txt', 'text/csv': '.csv',
            }
            ext = mime_map.get(content_type, '.bin')
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                return _convert_file(tmp_path)
            finally:
                os.unlink(tmp_path)
        except requests.RequestException as e:
            return f"Failed to download URL: {str(e)}"
        except Exception as e:
            return f"Error processing URL content: {str(e)}"
    else:
        return _convert_file(source)


def _convert_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    try:
        if ext == '.pdf':
            return _convert_pdf(path)
        elif ext == '.docx':
            return _convert_docx(path)
        elif ext == '.pptx':
            return _convert_pptx(path)
        elif ext == '.xlsx':
            return _convert_xlsx(path)
        elif ext in ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff'):
            return _convert_image(path)
        elif ext in ('.html', '.htm'):
            return _convert_html(path)
        elif ext == '.xml':
            return _convert_xml(path)
        elif ext == '.txt':
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.csv':
            return _convert_csv(path)
        else:
            return f"Unsupported file type: {ext}"
    except Exception as e:
        return f"Error converting {path}: {str(e)}"


def _convert_pdf(path: str) -> str:
    import pdf_inspector
    result = pdf_inspector.process_pdf(path)

    # Fast path: TextBased or Mixed - use markdown directly
    if result.pdf_type in ("TextBased", "Mixed"):
        if result.markdown:
            return result.markdown

    # Fallback for Scanned, ImageBased, or when markdown is None: use OCR
    import fitz  # PyMuPDF (lazy, needed for rendering pages to images)
    import pytesseract

    doc = fitz.open(path)
    texts = []
    for page in doc:
        pix = page.get_pixmap(dpi=300)
        img = pix.tobytes("png")
        text = pytesseract.image_to_string(img)
        texts.append(text)
    doc.close()
    return "\n".join(texts)


def _convert_docx(path: str) -> str:
    import docx  # python-docx (lazy)
    document = docx.Document(path)
    return "\n".join(para.text for para in document.paragraphs)


def _convert_pptx(path: str) -> str:
    from pptx import Presentation  # python-pptx (lazy)
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides.append(shape.text)
    return "\n".join(slides)


def _convert_xlsx(path: str) -> str:
    import openpyxl  # (lazy)
    import tabulate  # (lazy)
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets_md = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(c) if c is not None else "" for c in rows[0]]
        table = [[str(c) if c is not None else "" for c in row] for row in rows[1:]]
        md = tabulate.tabulate(table, headers=headers, tablefmt="github")
        sheets_md.append(f"### Sheet: {ws.title}\n\n{md}")
    wb.close()
    return "\n\n".join(sheets_md)


def _convert_image(path: str) -> str:
    from PIL import Image  # Pillow (lazy)
    import pytesseract  # (lazy)
    image = Image.open(path)
    if image.mode != "RGB":
        image = image.convert("RGB")
    return pytesseract.image_to_string(image)


def _convert_html(path: str) -> str:
    from bs4 import BeautifulSoup  # (lazy)
    import html2text  # (lazy)
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'html.parser')
    converter = html2text.HTML2Text()
    converter.ignore_links = False
    return converter.handle(str(soup))


def _convert_xml(path: str) -> str:
    from bs4 import BeautifulSoup  # (lazy)
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'xml')
    return soup.get_text(separator='\n', strip=True)


def _convert_csv(path: str) -> str:
    import csv  # (lazy)
    import tabulate  # (lazy)
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return ""
    return tabulate.tabulate(rows[1:], headers=rows[0], tablefmt="github")


@mcp.tool()
def read_code_outline(file_path: str) -> str:
    """Returns ONLY the function and class signatures of a Python file. Use this BEFORE reading full files to save tokens."""
    if len(file_path) > _MAX_PATH_LEN:
        return f"Error: path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(file_path)
    except ValueError as e:
        return f"Error: {e}"

    if _is_sensitive_file(resolved):
        return "Error: Access denied — sensitive file type."

    try:
        with open(resolved, "r") as f:
            source = f.read()
        tree = ast.parse(source)
    except Exception as e:
        return f"Failed to parse outline: {str(e)}"

    # Use a NodeVisitor to preserve class context for methods.
    class OutlineVisitor(ast.NodeVisitor):
        def __init__(self):
            self.outline = []
            self._class_stack = []

        def visit_ClassDef(self, node):
            self.outline.append(f"Class: {node.name}")
            self._class_stack.append(node.name)
            self.generic_visit(node)
            self._class_stack.pop()

        def visit_FunctionDef(self, node):
            args = [a.arg for a in node.args.args]
            if self._class_stack:
                self.outline.append(
                    f"  Method: {self._class_stack[-1]}.{node.name}({', '.join(args)})"
                )
            else:
                self.outline.append(f"Function: {node.name}({', '.join(args)})")
            self.generic_visit(node)

        visit_AsyncFunctionDef = visit_FunctionDef

    visitor = OutlineVisitor()
    visitor.visit(tree)
    return "\n".join(visitor.outline) if visitor.outline else "No classes or functions found."


@mcp.tool()
def run_command_compressed(command: str) -> str:
    """
    Runs a terminal command but truncates successful output to save tokens.
    The command is executed without a shell (shell=False) for security.
    Use standard shell syntax (pipes, redirects) is NOT supported.
    Preserves errors.
    """
    if len(command) > _MAX_CMD_LEN:
        return f"Error: command exceeds maximum length of {_MAX_CMD_LEN} characters."
    try:
        # Split safely — no shell=True to prevent injection.
        args = shlex.split(command)
    except ValueError as e:
        return f"Error: could not parse command: {e}"

    try:
        result = subprocess.run(
            args,
            shell=False,
            text=True,
            capture_output=True,
            timeout=60,
        )

        if result.returncode != 0:
            return f"FAILED. Error Trace:\n{result.stderr[-2000:]}"

        output = result.stdout.strip()
        if len(output) > 500:
            return f"SUCCESS. (Output truncated to save tokens): {output[:500]}..."
        return f"SUCCESS: {output}"

    except subprocess.TimeoutExpired:
        return "Error: command timed out after 60 seconds."
    except Exception as e:
        return f"Execution failed: {str(e)}"


@mcp.tool()
def compress_and_read_image(image_path: str) -> str:
    """Resizes and compresses large UI screenshots before analysis to save vision tokens."""
    if len(image_path) > _MAX_PATH_LEN:
        return f"Error: path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(image_path)
    except ValueError as e:
        return f"Error: {e}"

    if _is_sensitive_file(resolved):
        return "Error: Access denied — sensitive file type."

    try:
        from PIL import Image
        with Image.open(resolved) as img:
            img = img.convert("RGB")
            img.thumbnail((800, 800))

            buffer = io.BytesIO()
            img.save(buffer, format="JPEG", quality=60)
            encoded = base64.b64encode(buffer.getvalue()).decode("utf-8")

            return f"data:image/jpeg;base64,{encoded}"
    except Exception as e:
        return f"Image compression failed: {str(e)}"


@mcp.tool()
def get_workspace_info() -> dict:
    """
    Returns information about the workspace configuration to help the AI
    understand path mappings between host and container.
    
    Returns:
        dict with keys:
        - workspace_root: Container path to workspace root (e.g., "/workspace")
        - workspace_host_path: Host path mounted to workspace (if known)
        - workspace_project: Current WORKSPACE_PROJECT restriction (if any)
        - subdirectories: List of immediate subdirectories under workspace_root
    """
    info = {
        "workspace_root": _WORKSPACE,
        "workspace_host_path": os.environ.get("WORKSPACE_PATH", "unknown (set WORKSPACE_PATH in .env)"),
        "workspace_project": _WORKSPACE_PROJECT if _WORKSPACE_PROJECT else None,
        "subdirectories": []
    }
    try:
        with os.scandir(_WORKSPACE) as entries:
            info["subdirectories"] = sorted([entry.name for entry in entries if entry.is_dir()])
    except Exception:
        info["subdirectories"] = []
    return info


@mcp.tool()
def map_repository(directory: str = None, max_depth: int = 3) -> str:
    if directory is None:
        directory = _WORKSPACE
    """Always use this instead of 'ls' or 'tree' to understand the project structure."""
    if len(directory) > _MAX_PATH_LEN:
        return f"Error: directory path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(directory)
    except ValueError as e:
        return f"Error: {e}"

    try:
        # Respect .gitignore if it exists
        gitignore_path = os.path.join(resolved, '.gitignore')
        ignore_spec = None
        if os.path.exists(gitignore_path):
            with open(gitignore_path, 'r') as f:
                ignore_spec = pathspec.PathSpec.from_lines('gitwildmatch', f)

        repo_map = []
        file_count = 0
        MAX_FILES = 2000  # Hard cap

        for root, dirs, files in os.walk(resolved):
            # Calculate current depth
            depth = root.replace(resolved, '').count(os.sep)
            if depth >= max_depth:
                dirs[:] = []  # Stop walking deeper
                continue

            # Filter out junk directories immediately
            dirs[:] = [d for d in dirs if d not in {".git", "node_modules", "venv", "__pycache__", "dist", "build"}]

            # Format output tightly
            indent = '  ' * depth
            folder_name = os.path.basename(root) or "ROOT"
            repo_map.append(f"{indent}📂 {folder_name}/")

            for file in files:
                rel_path = os.path.relpath(os.path.join(root, file), resolved)
                if ignore_spec and ignore_spec.match_file(rel_path):
                    continue

                # Skip heavy binary/lock files
                if file.endswith(('.lock', '.png', '.sqlite', '.min.js', '.pdf', '.jpg', '.jpeg', '.gif', '.svg', '.ico', '.woff', '.woff2')):
                    continue

                repo_map.append(f"{indent}  📄 {file}")
                file_count += 1
                if file_count >= MAX_FILES:
                    repo_map.append("... (truncated)")
                    output = "\n".join(repo_map)
                    return output[:8000] + "\n... (truncated for size)"

        # Join and truncate to ~2000 tokens
        final_output = "\n".join(repo_map)
        if len(final_output) > 8000:
            return final_output[:8000] + "\n... (truncated for size)"
        return final_output
    except Exception as e:
        return f"Mapping failed: {str(e)}"


@mcp.tool()
def smart_code_search(keyword: str, file_pattern: str = "*.*") -> str:
    """Use this instead of 'grep'. Searches code and returns the top matches with surrounding context lines."""
    if len(keyword) > _MAX_QUERY_LEN:
        return f"Error: query exceeds maximum length of {_MAX_QUERY_LEN} characters."
    try:
        # Use ripgrep if available (faster), fallback to grep
        import shutil
        rg_path = shutil.which('rg')
        workspace = _WORKSPACE
        # Ensure workspace has trailing slash for replacement
        workspace_prefix = workspace if workspace.endswith(os.sep) else workspace + os.sep

        if rg_path:
            # -C 2: 2 lines of context, -n: show line numbers, --no-heading: cleaner output
            cmd = f'rg -C 2 -n --no-heading --type-add "search:*" --include "{file_pattern}" "{keyword}" "{workspace}" | head -n 100'
        else:
            cmd = f'grep -rn -C 2 --include="{file_pattern}" "{keyword}" "{workspace}" 2>/dev/null | head -n 100'

        result = subprocess.run(cmd, shell=True, capture_output=True, text=True)

        if not result.stdout:
            return f"No matches found for '{keyword}'."

        output_lines = result.stdout.strip().split('\n')

        # Compress the output format to save tokens
        compressed_results = []
        for line in output_lines:
            if line == "--":  # Grep's context separator
                compressed_results.append("---")
            else:
                # Remove the workspace prefix for brevity
                compressed_results.append(line.replace(workspace_prefix, ''))

        # Hard limit the return payload to ~2000 tokens
        final_text = "\n".join(compressed_results)
        if len(final_text) > 8000:
            final_text = final_text[:8000] + "\n... (truncated)"

        return f"Found matches (showing top results):\n{final_text}"

    except Exception as e:
        return f"Search failed: {str(e)}"


@mcp.tool()
def read_file_skeleton(file_path: str) -> str:
    """Always use this BEFORE reading a full file. Returns only the imports, functions, classes, and types."""
    if len(file_path) > _MAX_PATH_LEN:
        return f"Error: path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(file_path)
    except ValueError as e:
        return f"Error: {e}"

    if _is_sensitive_file(resolved):
        return "Error: Access denied — sensitive file type."

    try:
        with open(resolved, 'r') as f:
            lines = f.readlines()

        skeleton = []
        # Basic regex to catch Python, JS/TS, and Go signatures
        patterns = [
            r"^(def |class |import |from |@)",  # Python + decorators
            r"^(export |function |class |const .*=.*=>|type |interface |var |let |static )",  # JS/TS
            r"^(func |type |package |import |struct )"  # Go
        ]
        combined_pattern = re.compile("|".join(patterns))

        for i, line in enumerate(lines):
            stripped = line.lstrip()
            if combined_pattern.match(stripped):
                # Include the line number so the LLM knows exactly where to look later
                skeleton.append(f"Line {i+1}: {line.rstrip()}")

        if not skeleton:
            return "No obvious functions/classes/imports found. The file might be purely data or UI markup."

        # Limit skeleton to ~500 tokens to prevent abuse
        skeleton_text = "\n".join(skeleton)
        if len(skeleton_text) > 2000:
            skeleton_text = skeleton_text[:2000] + "\n... (truncated)"

        return "File Skeleton:\n" + skeleton_text
    except Exception as e:
        return f"Failed to read skeleton: {str(e)}"


@mcp.tool()
def focused_glob(pattern: str, directory: str = None, limit: int = 50) -> str:
    if directory is None:
        directory = _WORKSPACE
    """Finds files matching a pattern. Automatically filters junk and caps results to save tokens."""
    try:
        if len(directory) > _MAX_PATH_LEN:
            return f"Error: directory path exceeds maximum allowed length."
        try:
            resolved = _validate_local_path(directory)
        except ValueError as e:
            return f"Error: {e}"

        # Load .gitignore rules
        gitignore_path = os.path.join(resolved, '.gitignore')
        spec = None
        if os.path.exists(gitignore_path):
            with open(gitignore_path, 'r') as f:
                spec = pathspec.PathSpec.from_lines('gitwildmatch', f)

        # Recursive globbing - handle patterns with ** correctly
        full_pattern = os.path.join(resolved, pattern)
        raw_matches = glob.glob(full_pattern, recursive=True)

        filtered_matches = []
        for match in raw_matches:
            rel_path = os.path.relpath(match, resolved)

            # Filter by .gitignore
            if spec and spec.match_file(rel_path):
                continue

            # Filter common junk directories
            parts = rel_path.split(os.sep)
            if any(part in parts for part in {'.git', 'node_modules', 'dist', 'build', 'venv', '__pycache__'}):
                continue

            # Skip directories
            if os.path.isdir(match):
                continue

            filtered_matches.append(rel_path)
            if len(filtered_matches) >= limit:
                break

        if not filtered_matches:
            return "No matches found (ignoring hidden/junk folders)."

        result = "\n".join(filtered_matches)
        if len(raw_matches) > limit:
            result += f"\n\n... and {len(raw_matches) - limit} more files omitted. Be more specific with your pattern."

        return f"Matches found:\n{result}"
    except Exception as e:
        return f"Glob failed: {str(e)}"


@mcp.tool()
def read_lines(file_path: str, start_line: int, end_line: int) -> str:
    """Reads a specific range of lines from a file. Use this after read_file_skeleton to extract only the code you need."""
    if len(file_path) > _MAX_PATH_LEN:
        return f"Error: path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(file_path)
    except ValueError as e:
        return f"Error: {e}"

    if _is_sensitive_file(resolved):
        return "Error: Access denied — sensitive file type."

    try:
        if not os.path.exists(resolved):
            return f"Error: File {resolved} not found."

        if start_line < 1 or end_line < start_line:
            return "Error: Invalid line range."

        with open(resolved, 'r', encoding='utf-8', errors='ignore') as f:
            output = []
            for i, line in enumerate(f, 1):
                if i >= start_line and i <= end_line:
                    output.append(f"{i}: {line.rstrip()}")
                if i > end_line:
                    break

        if not output:
            return f"No lines found in range {start_line}-{end_line}."

        content = "\n".join(output)
        return f"File: {resolved} (Lines {start_line}-{end_line})\n---\n{content}"
    except Exception as e:
        return f"Failed to read lines: {str(e)}"


# ---------------------------------------------------------------------------
# FTS5 full-text search
# ---------------------------------------------------------------------------

def _init_fts_db(db_path: str = "/tmp/code_fts.db") -> str:
    """
    Initialize SQLite FTS5 database and index the workspace.
    - Clears stale rows before re-indexing to prevent duplicates on restart.
    - Only indexes files with allowlisted extensions.
    - Skips sensitive files.
    """
    workspace = _WORKSPACE_PROJECT if _WORKSPACE_PROJECT else _WORKSPACE
    workspace_prefix = workspace if workspace.endswith(os.sep) else workspace + os.sep
    ignore_dirs = {".git", "node_modules", "venv", "__pycache__", ".venv"}

    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()

    cursor.execute("""
        CREATE VIRTUAL TABLE IF NOT EXISTS code_fts
        USING fts5(
            path,
            content,
            line_num,
            tokenize='porter unicode61'
        )
    """)

    # Clear stale data before re-indexing to avoid duplicates on restart.
    cursor.execute("DELETE FROM code_fts")

    for root, dirs, files in os.walk(workspace):
        dirs[:] = [d for d in dirs if d not in ignore_dirs]
        for file in files:
            # Only index allowlisted file extensions.
            _, ext = os.path.splitext(file)
            if ext.lower() not in _INDEXABLE_EXTENSIONS:
                continue
            # Skip sensitive files.
            if _is_sensitive_file(file):
                continue

            filepath = os.path.join(root, file)
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    for i, line in enumerate(f, 1):
                        content = line.strip()
                        if content:
                            cursor.execute(
                                "INSERT INTO code_fts (path, content, line_num) VALUES (?, ?, ?)",
                                (filepath.replace(workspace, "").lstrip("/") or filepath, content, i)
                            )
            except (OSError, IOError, UnicodeDecodeError):
                continue

    conn.commit()
    conn.close()
    return db_path


# Run DB initialization in a background thread so it doesn't block startup.
FTS_DB: str | None = None
_fts_lock = threading.Lock()


def _fts_init_worker():
    global FTS_DB
    try:
        db = _init_fts_db()
        with _fts_lock:
            FTS_DB = db
    except Exception as e:
        print(f"Warning: FTS5 initialization failed: {e}")


threading.Thread(target=_fts_init_worker, daemon=True).start()


@mcp.tool()
def search_codebase(query: str, limit: int = 5) -> str:
    """
    Full-text search through codebase using SQLite FTS5 (BM25 ranking).
    Returns top matches with file paths and line numbers.
    Much more precise than grep and token-efficient.
    """
    if len(query) > _MAX_QUERY_LEN:
        return f"Error: query exceeds maximum length of {_MAX_QUERY_LEN} characters."

    with _fts_lock:
        db = FTS_DB

    if not db:
        return "FTS5 search not available. Database may still be initializing — try again shortly."

    try:
        conn = sqlite3.connect(db)
        cursor = conn.cursor()

        cursor.execute("""
            SELECT path, line_num, content
            FROM code_fts
            WHERE code_fts MATCH ?
            ORDER BY bm25(code_fts)
            LIMIT ?
        """, (query, limit))

        results = cursor.fetchall()
        conn.close()

        if not results:
            return f"No matches found for '{query}'"

        formatted = [f"{path}:{line_num}: {content}" for path, line_num, content in results]
        return "\n".join(formatted)
    except Exception as e:
        return f"Search failed: {str(e)}"


@mcp.tool()
def safe_read_file(file_path: str, force: bool = False) -> str:
    """
    Always use this to read code. It prevents accidentally reading massive files.
    Checks file size first to avoid token bloat from large legacy files.
    Only files within the configured workspace directory are accessible.
    """
    if len(file_path) > _MAX_PATH_LEN:
        return f"Error: path exceeds maximum allowed length."
    try:
        resolved = _validate_local_path(file_path)
    except ValueError as e:
        return f"Error: {e}"

    if _is_sensitive_file(resolved):
        return "Error: Access denied — sensitive file type."

    try:
        if not os.path.exists(resolved):
            return f"Error: File '{resolved}' does not exist."

        size_bytes = os.path.getsize(resolved)
        # 30 KB is roughly 8,000 tokens. Adjust threshold as needed.
        if size_bytes > 30000 and not force:
            return (
                f"BLOCKED: {resolved} is too large ({size_bytes / 1000:.1f} KB). "
                f"To save tokens, do not read this entire file. "
                f"Use search_codebase() to find the specific lines you need, "
                f"or call this tool again with force=True if absolutely necessary."
            )

        with open(resolved, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Error reading file: {str(e)}"


# ---------------------------------------------------------------------------
# SearXNG-to-4get Gateway (embedded HTTP server)
# ---------------------------------------------------------------------------
# Starts a FastAPI gateway on port 8080 inside the MCP server container.
# CRW connects here thinking it's talking to SearXNG; we translate to 4get.

def _start_gateway():
    """Run the gateway FastAPI server in a daemon thread."""
    import uvicorn
    from gateway import gateway_app, FOURGET_URL as _  # noqa: F811

    # Set FOURGET_URL in the gateway module
    import gateway
    gateway.FOURGET_URL = os.getenv("FOURGET_URL", "http://localhost:80")

    logging.getLogger("gateway").info(
        "Starting SearXNG-to-4get gateway on 0.0.0.0:8080 (backend: %s)",
        gateway.FOURGET_URL,
    )
    config = uvicorn.Config(
        gateway.gateway_app,
        host="0.0.0.0",
        port=8080,
        log_level="warning",
        access_log=False,
    )
    server = uvicorn.Server(config)
    server.run()


# Start gateway in a background daemon thread so it doesn't block MCP stdio
_gateway_thread = threading.Thread(target=_start_gateway, daemon=True, name="gateway")
_gateway_thread.start()


if __name__ == "__main__":
    mcp.run()
